import os
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import List, Optional
from pptx import Presentation
import google.generativeai as genai
import json

app = FastAPI()

# Configure Gemini
genai.configure(api_key=os.environ.get("GOOGLE_API_KEY"))
model = genai.GenerativeModel('gemini-2.5-flash')

# ── Pydantic Models ──────────────────────────────────────────────────────────

class RAGRequest(BaseModel):
    """Request model for RAG queries"""
    query: str
    documents: List[str] = []
    top_k: int = 5

class RAGSourceDocument(BaseModel):
    """Source document reference"""
    title: str
    content: str
    score: float
    knowledge_base_id: int = 0

class RAGResponse(BaseModel):
    """Response model for RAG queries"""
    answer: str
    sources: List[RAGSourceDocument] = []
    confidence: float = 0.8

class QuizQuestion(BaseModel):
    """Single quiz question"""
    question: str
    options: List[str]
    correct_answer: str

class QuizResponse(BaseModel):
    """Quiz generation response"""
    quiz: List[QuizQuestion]

# ── Helper Functions ────────────────────────────────────────────────────────

def extract_text_from_pptx(file_path: str) -> str:
    """Extract text from PPTX file"""
    prs = Presentation(file_path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
    return "\n".join(full_text)

def simple_relevance_score(query: str, document: str) -> float:
    """
    Simple relevance scoring based on keyword overlap.
    Enhanced to give higher scores for healthcare/billing documents.
    """
    query_words = set(query.lower().split())
    doc_words = set(document.lower().split())
    
    if len(doc_words) == 0:
        return 0.5  # Default score for empty documents
    
    # Calculate overlap
    overlap = len(query_words & doc_words)
    
    # Boost score if document contains common healthcare/billing keywords
    healthcare_keywords = {'medical', 'billing', 'insurance', 'patient', 'copay', 'deductible', 
                          'coinsurance', 'claim', 'coverage', 'healthcare', 'hospital', 
                          'appointment', 'provider', 'treatment', 'procedure'}
    
    has_healthcare_content = bool(doc_words & healthcare_keywords)
    
    # Base score on overlap
    base_score = overlap / max(len(query_words), 1)
    
    # Boost score if it's healthcare-related content
    if has_healthcare_content:
        base_score = max(base_score, 0.6)  # Minimum 60% for healthcare docs
    
    # Ensure score is between 0.5 and 1.0
    return min(max(base_score, 0.5), 1.0)

def extract_document_title(document: str) -> str:
    """
    Extract the document title from the formatted content.
    Looks for 'DOCUMENT: [title]' or 'FILE: [filename]' patterns.
    """
    lines = document.split('\n')
    for line in lines[:10]:  # Check first 10 lines
        if line.startswith('DOCUMENT:'):
            return line.replace('DOCUMENT:', '').strip()
        if line.startswith('FILE:'):
            filename = line.replace('FILE:', '').strip()
            # Remove file extension for cleaner title
            return filename.rsplit('.', 1)[0] if '.' in filename else filename
    return "Healthcare Document"

def rank_documents_by_relevance(query: str, documents: List[str], top_k: int = 5) -> List[RAGSourceDocument]:
    """
    Rank documents by relevance to query and return top_k.
    """
    ranked = []
    for i, doc in enumerate(documents):
        score = simple_relevance_score(query, doc)
        
        # Extract proper title from document
        title = extract_document_title(doc)
        
        ranked.append(RAGSourceDocument(
            title=title,
            content=doc[:500],  # Truncate for response
            score=max(0.5, score),  # Minimum score of 0.5 for all documents
            knowledge_base_id=i
        ))
    
    # Sort by score descending
    ranked.sort(key=lambda x: x.score, reverse=True)
    return ranked[:top_k]

# ── RAG Endpoints ───────────────────────────────────────────────────────────

@app.post("/api/rag/query", response_model=RAGResponse)
async def rag_query(request: RAGRequest):
    """
    Query RAG service with documents.
    Returns AI-generated answer based on provided documents + referenced sources.
    """
    try:
        if not request.query or not request.query.strip():
            raise HTTPException(status_code=400, detail="Query cannot be empty")
        
        # Rank documents by relevance
        relevant_docs = rank_documents_by_relevance(
            request.query, 
            request.documents, 
            request.top_k
        )
        
        # Prepare context for Gemini
        context = "\n\n---\n\n".join([doc.content for doc in relevant_docs])
        
        # Construct prompt for Gemini
        prompt = f"""You are a helpful healthcare administrative assistant with expertise in medical billing, insurance, and healthcare procedures.

You have been provided with knowledge base documents that may contain relevant information to answer the user's question.

KNOWLEDGE BASE DOCUMENTS:
{context}

USER QUESTION:
{request.query}

INSTRUCTIONS:
- Provide a clear, helpful, and professional answer to the user's question
- Use information from the knowledge base documents when available
- If the documents contain relevant information about the topic, use it to provide a comprehensive answer
- You can explain healthcare administrative concepts, billing terminology, and procedures based on the context provided
- Keep responses conversational and easy to understand (2-4 sentences)
- Focus on administrative and billing topics, not clinical medical advice
- Be helpful and informative - don't say "I don't have information" unless the question is completely unrelated to healthcare administration

Based on the knowledge base provided above, please answer the user's question in a helpful and informative way:"""
        
        # Call Gemini API
        response = model.generate_content(prompt)
        answer = response.text
        
        # Calculate overall confidence based on average document scores
        avg_confidence = sum(doc.score for doc in relevant_docs) / len(relevant_docs) if relevant_docs else 0.5
        
        return RAGResponse(
            answer=answer,
            sources=relevant_docs,
            confidence=min(avg_confidence, 1.0)
        )
    
    except HTTPException as e:
        raise e
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

# ── Quiz Generation Endpoint (Original) ─────────────────────────────────────

@app.post("/generate-quiz", response_model=QuizResponse)
async def generate_quiz(file: UploadFile = File(...)):
    """
    Generate a 20-question multiple choice quiz from a PowerPoint file.
    """
    if not file.filename.endswith('.pptx'):
        raise HTTPException(status_code=400, detail="Only .pptx files are supported.")
    
    # Save temp file to parse
    temp_path = f"temp_{file.filename}"
    with open(temp_path, "wb") as buffer:
        buffer.write(await file.read())
    
    try:
        # 1. Parse PPTX
        context_text = extract_text_from_pptx(temp_path)
        
        # 2. Construct Prompt
        prompt = f"""Extract the most important information from the following text and create a quiz.
Requirements:
- Exactly 20 Multiple Choice Questions
- 4 options per question
- Clearly indicate the correct answer
- Return the result strictly as a JSON array with this structure:
[
  {{
    "question": "question text",
    "options": ["option1", "option2", "option3", "option4"],
    "correct_answer": "option1"
  }},
  ...
]

Text: {context_text[:10000]}"""
        
        # 3. Call Gemini
        response = model.generate_content(prompt)
        
        # Clean up the AI response (removing markdown code blocks if present)
        raw_text = response.text.replace('```json', '').replace('```', '').strip()
        quiz_data = json.loads(raw_text)
        
        return {"quiz": quiz_data}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)

# ── Health Check ────────────────────────────────────────────────────────────

@app.get("/")
def health_check():
    """Health check endpoint"""
    return {"status": "Service is running"}

@app.get("/api/health")
def api_health_check():
    """API health check endpoint"""
    return {"status": "running", "version": "1.0.0"}
